from flask import Flask, render_template, request, redirect, url_for, session, flash, send_file, Response
from flask_sqlalchemy import SQLAlchemy
from flask_wtf import FlaskForm
from wtforms import StringField, SubmitField, PasswordField, TextAreaField
from wtforms.validators import DataRequired, Length, Email
from nltk.tokenize import sent_tokenize
from collections import Counter
from bs4 import BeautifulSoup
import requests
import os
import pyttsx3
from pptx import Presentation
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
from io import BytesIO
import tempfile
import pytesseract
from PIL import Image
import nltk
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
from reportlab.lib.utils import simpleSplit
from io import BytesIO
import nltk
nltk.download('punkt')
from nltk.tokenize import sent_tokenize

# Download the standard 'punkt' package (if not already downloaded)
nltk.download('punkt')

app = Flask(__name__)
app.secret_key = 'your_secret_key'

# Specify the path to the Tesseract executable
pytesseract.pytesseract.tesseract_cmd = '/usr/bin/tesseract'  # Update this path as needed

app.config['SECRET_KEY'] = 'your_secret_key'
app.config['SQLALCHEMY_DATABASE_URI'] = 'sqlite:///users.db'
app.config['SQLALCHEMY_TRACK_MODIFICATIONS'] = False
app.config['UPLOADED_PHOTOS_DEST'] = 'uploads'

# Ensure the upload directory exists
if not os.path.exists(app.config['UPLOADED_PHOTOS_DEST']):
    os.makedirs(app.config['UPLOADED_PHOTOS_DEST'])

db = SQLAlchemy(app)

class User(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    username = db.Column(db.String(50), unique=True, nullable=False)
    password = db.Column(db.String(100), nullable=False)
    email = db.Column(db.String(100), unique=True, nullable=False)
    fullname = db.Column(db.String(100), nullable=True)
    education = db.Column(db.String(100), nullable=True)

class RegistrationForm(FlaskForm):
    username = StringField('Username', validators=[DataRequired(), Length(min=4, max=25)])
    email = StringField('Email', validators=[DataRequired(), Email()])
    password = PasswordField('Password', validators=[DataRequired(), Length(min=6)])
    submit = SubmitField('Sign Up')

class LoginForm(FlaskForm):
    username = StringField('Username', validators=[DataRequired()])
    password = PasswordField('Password', validators=[DataRequired()])
    submit = SubmitField('Login')

@app.route('/')
def index():
    if 'username' in session:
        return redirect(url_for('dashboard'))
    return render_template('index.html')

@app.route('/register', methods=['GET', 'POST'])
def register():
    form = RegistrationForm()
    if form.validate_on_submit():
        new_user = User(username=form.username.data, email=form.email.data, password=form.password.data)
        db.session.add(new_user)
        db.session.commit()
        flash('Registration successful! Please log in.', 'success')
        return redirect(url_for('login'))
    return render_template('register.html', form=form)

@app.route('/login', methods=['GET', 'POST'])
def login():
    form = LoginForm()
    if form.validate_on_submit():
        user = User.query.filter_by(username=form.username.data, password=form.password.data).first()
        if user:
            session['username'] = user.username
            flash('Login successful!', 'success')
            return redirect(url_for('dashboard'))
        else:
            flash('Invalid username or password', 'danger')
    return render_template('login.html', form=form)

@app.route('/dashboard')
def dashboard():
    if 'username' not in session:
        return redirect(url_for('login'))
    return render_template('dashboard.html')

@app.route('/logout')
def logout():
    session.pop('username', None)
    flash('You have been logged out.', 'success')
    return redirect(url_for('login'))

@app.route('/image-summarization', methods=['GET', 'POST'])
def image_summarization():
    if 'username' not in session:
        return redirect(url_for('login'))
    if request.method == 'POST':
        if 'photo' not in request.files:
            flash('No file part', 'danger')
            return redirect(request.url)
        file = request.files['photo']
        if file.filename == '':
            flash('No selected file', 'danger')
            return redirect(request.url)
        if file:
            file_path = os.path.join(app.config['UPLOADED_PHOTOS_DEST'], file.filename)
            file.save(file_path)
            text = extract_text_from_image(file_path)
            num_sentences = int(request.form['num_sentences'])
            summary = summarize_text(text, num_sentences)
            session['summary'] = summary
            return redirect(url_for('summary'))
    return render_template('image_summarization.html')

@app.route('/text-summarization', methods=['GET', 'POST'])
def text_summarization():
    if request.method == 'POST':
        text = request.form['text']
        summary = generate_summary(text)  # Generate summary
        session['summary'] = summary  # Store in session
        return redirect(url_for('summary'))
    
    # If it's a GET request, just render the form page
    return render_template('text_summarization.html')


@app.route('/summary')
def summary():
    if 'summary' not in session:
        return redirect(url_for('dashboard'))
    return render_template('summary.html', summary=session['summary'])

@app.route('/enter-url', methods=['GET', 'POST'])
def enter_url():
    if request.method == 'POST':
        url = request.form['url']
        return redirect(url_for('summarize_website', url=url))
    return render_template('enter_url.html')

@app.route('/summarize-website', methods=['GET', 'POST'])
def summarize_website():
    if request.method == 'POST':
        url = request.form['url']
        num_sentences = int(request.form['num_sentences'])
        try:
            response = requests.get(url)
            response.raise_for_status()
            html_content = response.text
            soup = BeautifulSoup(html_content, 'html.parser')
            text_content = ' '.join(soup.stripped_strings)
            summary = summarize_text(text_content, num_sentences)
            sentences = sent_tokenize(summary)
            return render_template('website_summary.html', sentences=sentences)
        except Exception as e:
            error_message = f"An error occurred: {e}"
            flash(error_message, 'danger')
            return redirect(url_for('index'))
    return render_template('summarize_website.html')

@app.route('/download-pdf', methods=['POST'])
def download_pdf():
    summary_text = session.get('summary', '')  # Fetch summary from session

    if not summary_text:
        flash("No summary available to download.", "danger")
        return redirect(url_for('summary'))

    # Create a PDF in memory
    pdf_buffer = BytesIO()
    pdf = canvas.Canvas(pdf_buffer, pagesize=letter)
    pdf.setFont("Helvetica", 12)

    # Title
    pdf.drawString(100, 750, "Summary")

    # Set the starting position for the text
    text_object = pdf.beginText(100, 730)
    text_object.setFont("Helvetica", 12)

    # Define maximum width for text (adjust as needed)
    max_width = 400

    # Split text into lines that fit within the width
    lines = []
    for line in summary_text.split("\n"):
        wrapped_lines = simpleSplit(line, "Helvetica", 12, max_width)
        lines.extend(wrapped_lines)

    # Adjust vertical positioning and add text
    y_position = 730
    for line in lines:
        if y_position <= 50:  # Create a new page if reaching bottom
            pdf.showPage()
            pdf.setFont("Helvetica", 12)
            y_position = 750  # Reset position for new page

        pdf.drawString(100, y_position, line)
        y_position -= 20  # Move to the next line

    pdf.showPage()
    pdf.save()

    pdf_buffer.seek(0)

    return send_file(pdf_buffer, as_attachment=True, download_name='summary.pdf', mimetype='application/pdf')


@app.route('/download-pptx', methods=['POST'])
def download_pptx():
    summary_text = session.get('summary', '')  # Fetch from session

    if not summary_text:
        flash("No summary available to download.", "danger")
        return redirect(url_for('summary'))

    prs = Presentation()
    slide_layout = prs.slide_layouts[1]  # Title + Content Layout
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]

    title.text = "Summary"
    content.text = summary_text  # Insert summary text

    temp = tempfile.NamedTemporaryFile(delete=False, suffix='.pptx')
    prs.save(temp.name)
    temp.seek(0)

    return send_file(temp.name, as_attachment=True, download_name='summary.pptx')

@app.route('/text-to-speech', methods=['POST'])
def text_to_speech():
    if 'summary' in session:
        summary_text = session['summary']
        engine = pyttsx3.init()
        engine.say(summary_text)
        engine.runAndWait()
        return redirect(url_for('summary'))
    else:
        return redirect(url_for('dashboard'))

def extract_text_from_image(image_path):
    text = pytesseract.image_to_string(Image.open(image_path))
    return text

def ensure_punkt_tab():
    """
    Ensure the NLTK resource 'punkt_tab' is available.
    """
    try:
        nltk.data.find('tokenizers/punkt_tab')
    except LookupError:
        nltk.download('punkt_tab')

def summarize_text(text, num_sentences):
    # Ensure the required NLTK resource is available
    ensure_punkt_tab()
    sentences = sent_tokenize(text)
    if len(sentences) <= num_sentences:
        return text
    word_counts = Counter(word for sentence in sentences for word in sentence.split())
    sentence_scores = [(sum(word_counts[word] for word in sentence.split()), sentence) for sentence in sentences]
    top_sentences = sorted(sentence_scores, reverse=True)[:num_sentences]
    top_sentences.sort(key=lambda x: sentences.index(x[1]))
    summary = ' '.join(sentence for score, sentence in top_sentences)
    return summary

import nltk
from nltk.tokenize import sent_tokenize

# Download nltk data (only once)
nltk.download('punkt')

def generate_summary(text):
    """Generates a summarized version of the input text"""
    sentences = sent_tokenize(text)  # Splitting text into sentences
    summary = "\n".join(sentences[:5])  # Taking first 5 sentences as summary
    return summary


if __name__ == '__main__':
    with app.app_context():
        db.create_all()
    app.run(debug=True)
